"""Generate COFINFAD Analytics Poster as a PowerPoint file — WHITE background, teal accents."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Color palette (white bg + teal accents, matching reference poster) ---
TEAL = RGBColor(0x1A, 0x8F, 0xAA)
TEAL_DARK = RGBColor(0x14, 0x6E, 0x82)
TEAL_LIGHT_BG = RGBColor(0xE8, 0xF4, 0xF7)  # very light teal for card bg
NAVY = RGBColor(0x0C, 0x1A, 0x30)
NAVY_HEADER = RGBColor(0x0D, 0x1F, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)  # near-black for body text
MEDIUM_TEXT = RGBColor(0x47, 0x55, 0x69)  # medium gray for secondary text
MUTED_TEXT = RGBColor(0x64, 0x74, 0x8B)  # muted for labels
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)  # light gray card background
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)  # subtle border
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
SMU_BLUE = RGBColor(0x00, 0x3D, 0x7C)
POSTER_BORDER = RGBColor(0x1A, 0x8F, 0xAA)

# Slide dimensions: widescreen 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# --- Helper functions ---
def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.02
    return shape

def add_box(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox

def set_text(tf, text, size=10, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"
    return p

def add_para(tf, text, size=10, color=DARK_TEXT, bold=False, before=Pt(2), after=Pt(2), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = before
    p.space_after = after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"
    return p

def add_mixed(tf, segments, size=10, before=Pt(2), after=Pt(2), align=PP_ALIGN.LEFT):
    """segments: list of (text, color, bold)"""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = before
    p.space_after = after
    for text, color, bold in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Calibri"
    return p

def section_title(tf, title, size=11):
    p = tf.paragraphs[0] if tf.paragraphs and tf.paragraphs[0].text == "" else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = title.upper()
    run.font.size = Pt(size)
    run.font.color.rgb = TEAL
    run.font.bold = True
    run.font.name = "Calibri"
    return p

# ============================
# WHITE BACKGROUND
# ============================
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()
bg.shadow.inherit = False

# Outer poster border (teal)
outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.1), Inches(0.1),
                                prs.slide_width - Inches(0.2), prs.slide_height - Inches(0.2))
outer.fill.background()
outer.line.color.rgb = TEAL
outer.line.width = Pt(3)
outer.shadow.inherit = False

# ============================
# HEADER BAR (dark navy)
# ============================
header_h = Inches(0.9)
header = add_box(slide, Inches(0.1), Inches(0.1), prs.slide_width - Inches(0.2), header_h, NAVY_HEADER)

# Teal accent line under header
add_box(slide, Inches(0.1), Inches(0.1) + header_h - Pt(4), prs.slide_width - Inches(0.2), Pt(4), TEAL)

# SMU Logo
smu_box = add_rect(slide, Inches(0.35), Inches(0.25), Inches(0.65), Inches(0.5), SMU_BLUE)
smu_tb = add_text_box(slide, Inches(0.35), Inches(0.27), Inches(0.65), Inches(0.48))
set_text(smu_tb.text_frame, "SMU", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# SMU info text
smu_info = add_text_box(slide, Inches(1.1), Inches(0.22), Inches(2.8), Inches(0.6))
set_text(smu_info.text_frame, "SINGAPORE MANAGEMENT UNIVERSITY", size=7, color=RGBColor(0xA0,0xB4,0xCE), bold=True)
add_para(smu_info.text_frame, "School of Computing & Information Systems", size=7, color=RGBColor(0x8A,0xA3,0xC8), before=Pt(1), after=Pt(0))
add_para(smu_info.text_frame, "ISSS608 Visual Analytics & Applications", size=7, color=RGBColor(0x8A,0xA3,0xC8), before=Pt(1), after=Pt(0))

# Title (center)
title_tb = add_text_box(slide, Inches(4), Inches(0.12), Inches(5.5), Inches(0.8))
set_text(title_tb.text_frame, "COFINFAD Analytics", size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_para(title_tb.text_frame, "Customer Lifecycle & Transactional Behavior Analysis for Colombian Fintech",
         size=10, color=TEAL, bold=True, align=PP_ALIGN.CENTER, before=Pt(2))
add_para(title_tb.text_frame, "Ji Guofang  •  Lin Yan  •  Nguyen Trong Nhan   |   Team 13  •  April 2026",
         size=8, color=RGBColor(0x8A,0xA3,0xC8), align=PP_ALIGN.CENTER, before=Pt(2))

# COFINFAD brand (right)
brand_tb = add_text_box(slide, Inches(10.5), Inches(0.25), Inches(2.5), Inches(0.5))
set_text(brand_tb.text_frame, "COFINFAD", size=16, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
add_para(brand_tb.text_frame, "Visual Analytics Platform", size=7, color=RGBColor(0x8A,0xA3,0xC8), align=PP_ALIGN.RIGHT)

# ============================
# ROW 1: Motivation | Dataset | Methodology
# ============================
row1_top = Inches(1.15)
col_w = Inches(4.1)
col_gap = Inches(0.15)
margin_left = Inches(0.25)
card_h = Inches(2.55)

# --- Card 1: Motivation & Objectives ---
c1_left = margin_left
card1 = add_rect(slide, c1_left, row1_top, col_w, card_h, CARD_BG, CARD_BORDER)

tb1 = add_text_box(slide, c1_left + Inches(0.12), row1_top + Inches(0.08), col_w - Inches(0.24), card_h - Inches(0.12))
section_title(tb1.text_frame, "Motivation & Objectives")
add_para(tb1.text_frame, "Colombia's fintech sector is growing rapidly, yet many institutions lack accessible tools to analyze customer behavior and predict churn. This project aims to:",
         size=9, color=DARK_TEXT, after=Pt(4))

bullets = [
    ("Democratize analytics", " — enable non-technical stakeholders to explore customer data through interactive visualizations"),
    ("Identify behavioral segments", " — uncover distinct customer groups via unsupervised clustering on transactional features"),
    ("Predict churn risk", " — build interpretable models that quantify churn probability and surface key drivers"),
    ("Enable what-if analysis", " — let decision-makers simulate interventions and assess impact on retention"),
]
for title, desc in bullets:
    add_mixed(tb1.text_frame, [
        ("•  ", TEAL, True),
        (title, DARK_TEXT, True),
        (desc, MEDIUM_TEXT, False),
    ], size=8, before=Pt(1), after=Pt(1))

# Core question highlight
add_para(tb1.text_frame, "", size=3, before=Pt(2), after=Pt(0))
add_para(tb1.text_frame, "CORE QUESTION", size=8, color=TEAL, bold=True, before=Pt(3))
add_para(tb1.text_frame, '"How can interactive visual analytics help Colombian fintech institutions understand, segment, and retain their customers?"',
         size=8, color=DARK_TEXT, before=Pt(1))

# --- Card 2: Dataset Overview ---
c2_left = margin_left + col_w + col_gap
card2 = add_rect(slide, c2_left, row1_top, col_w, card_h, CARD_BG, CARD_BORDER)

tb2 = add_text_box(slide, c2_left + Inches(0.12), row1_top + Inches(0.08), col_w - Inches(0.24), card_h - Inches(0.12))
section_title(tb2.text_frame, "Dataset Overview")
add_para(tb2.text_frame, "The COFINFAD (Colombian Fintech Financial Analytics Dataset) contains comprehensive customer and transaction records:",
         size=9, color=DARK_TEXT, after=Pt(6))

# KPI boxes (teal background)
kpi_data = [("48,723", "Customers"), ("3.16M", "Transactions"), ("54", "Variables")]
kpi_w = Inches(1.18)
kpi_h = Inches(0.52)
kpi_gap = Inches(0.1)
kpi_top = row1_top + Inches(0.75)
for i, (val, label) in enumerate(kpi_data):
    kx = c2_left + Inches(0.15) + i * (kpi_w + kpi_gap)
    kpi_box = add_rect(slide, kx, kpi_top, kpi_w, kpi_h, TEAL_LIGHT_BG, TEAL, Pt(1))
    kpi_tb = add_text_box(slide, kx, kpi_top + Inches(0.03), kpi_w, kpi_h)
    set_text(kpi_tb.text_frame, val, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_para(kpi_tb.text_frame, label.upper(), size=7, color=MUTED_TEXT, align=PP_ALIGN.CENTER, before=Pt(0))

# Second KPI row
kpi_data2 = [("15", "Departments"), ("12 mo", "Time Span (2023)")]
kpi2_top = kpi_top + kpi_h + Inches(0.08)
kpi2_w = Inches(1.82)
for i, (val, label) in enumerate(kpi_data2):
    kx = c2_left + Inches(0.15) + i * (kpi2_w + kpi_gap)
    kpi_box = add_rect(slide, kx, kpi2_top, kpi2_w, kpi_h, TEAL_LIGHT_BG, TEAL, Pt(1))
    kpi_tb = add_text_box(slide, kx, kpi2_top + Inches(0.03), kpi2_w, kpi_h)
    set_text(kpi_tb.text_frame, val, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_para(kpi_tb.text_frame, label.upper(), size=7, color=MUTED_TEXT, align=PP_ALIGN.CENTER, before=Pt(0))

# Dataset descriptions
desc_top = kpi2_top + kpi_h + Inches(0.08)
desc_tb = add_text_box(slide, c2_left + Inches(0.12), desc_top, col_w - Inches(0.24), Inches(0.8))
set_text(desc_tb.text_frame, "", size=1)
add_mixed(desc_tb.text_frame, [
    ("Customer table: ", DARK_TEXT, True),
    ("demographics, income bracket, acquisition channel, churn probability, NPS, satisfaction score", MEDIUM_TEXT, False),
], size=8, before=Pt(0))
add_mixed(desc_tb.text_frame, [
    ("Transaction log: ", DARK_TEXT, True),
    ("date, amount (COP), type (Transfer, Withdrawal, Payment, Deposit), behavioral metrics", MEDIUM_TEXT, False),
], size=8, before=Pt(1))

# --- Card 3: Methodology & Approach ---
c3_left = margin_left + 2 * (col_w + col_gap)
card3 = add_rect(slide, c3_left, row1_top, col_w, card_h, CARD_BG, CARD_BORDER)

tb3 = add_text_box(slide, c3_left + Inches(0.12), row1_top + Inches(0.08), col_w - Inches(0.24), card_h - Inches(0.12))
section_title(tb3.text_frame, "Methodology & Approach")
add_para(tb3.text_frame, "We follow a 5-stage analytical framework that progresses from exploration to actionable simulation:",
         size=9, color=DARK_TEXT, after=Pt(4))

steps = [
    ("1", "EXPLORE", "EDA with interactive charts & maps"),
    ("2", "CONFIRM", "Hypothesis testing via ggstatsplot"),
    ("3", "SEGMENT", "K-Means, Hierarchical, PAM clustering"),
    ("4", "CALIBRATE", "Lasso logistic regression for churn"),
    ("5", "SIMULATE", "What-if analysis with real-time sliders"),
]
for num, name, desc in steps:
    add_mixed(tb3.text_frame, [
        (f"  {num}  ", WHITE, True),
        (f"  {name}", DARK_TEXT, True),
        (f"  —  {desc}", MUTED_TEXT, False),
    ], size=8, before=Pt(2), after=Pt(1))

# Teal badges for step numbers
step_top_base = row1_top + Inches(0.72)
for i in range(5):
    sy = step_top_base + Inches(i * 0.245)
    badge = add_rect(slide, c3_left + Inches(0.15), sy, Inches(0.22), Inches(0.18), TEAL)
    badge_tb = add_text_box(slide, c3_left + Inches(0.15), sy - Pt(1), Inches(0.22), Inches(0.18))
    set_text(badge_tb.text_frame, steps[i][0], size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_mixed(tb3.text_frame, [("Technology stack: ", DARK_TEXT, True)], size=9, before=Pt(6))
add_para(tb3.text_frame, "R Shiny  •  Plotly  •  ggplot2  •  ggstatsplot  •  Leaflet  •  glmnet  •  cluster  •  factoextra  •  bslib  •  DT",
         size=7, color=TEAL, bold=True, before=Pt(2))

# ============================
# ROW 2: Shiny Application - 5 Modules
# ============================
row2_top = row1_top + card_h + Inches(0.12)
row2_title_h = Inches(0.38)

# Section title
title_bar = add_text_box(slide, margin_left, row2_top, prs.slide_width - 2 * margin_left, row2_title_h)
set_text(title_bar.text_frame, "INTERACTIVE SHINY APPLICATION — 5 ANALYTICAL MODULES",
         size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_para(title_bar.text_frame, "All modules are integrated into a single dark-themed R Shiny dashboard at guofang1028.shinyapps.io/ShinyApp",
         size=8, color=MUTED_TEXT, align=PP_ALIGN.CENTER, before=Pt(1))

mod_top = row2_top + row2_title_h + Inches(0.05)
mod_w = Inches(2.45)
mod_h = Inches(1.4)
mod_gap = Inches(0.12)

modules = [
    ("MODULE 1", "Exploratory Data Analysis",
     "Interactive time-series charts, leaflet choropleth maps of regional distribution, and distribution plots by income bracket and transaction type."),
    ("MODULE 2", "Confirmatory Analysis",
     "Statistical hypothesis testing with ggbetweenstats (t-test, Welch, Mann-Whitney, Kruskal-Wallis) and correlation matrices."),
    ("MODULE 3", "Customer Clustering",
     "Multi-algorithm segmentation (K-Means, Hierarchical, PAM) with PCA scatter plots, silhouette validation, and heatmaps."),
    ("MODULE 4", "Churn Prediction",
     "Lasso (L1) logistic regression with cross-validation, variable importance ranking, and residual diagnostics."),
    ("MODULE 5", "What-If Simulator",
     "Interactive sliders to adjust customer attributes and observe real-time churn probability changes with sensitivity analysis."),
]

for i, (mod_num, mod_title, mod_desc) in enumerate(modules):
    mx = margin_left + i * (mod_w + mod_gap)
    mod_card = add_rect(slide, mx, mod_top, mod_w, mod_h, WHITE, CARD_BORDER)

    # Module number badge (teal)
    badge = add_rect(slide, mx + Inches(0.1), mod_top + Inches(0.08), Inches(0.7), Inches(0.18), TEAL)
    badge_tb = add_text_box(slide, mx + Inches(0.1), mod_top + Inches(0.06), Inches(0.7), Inches(0.2))
    set_text(badge_tb.text_frame, mod_num, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    mod_tb = add_text_box(slide, mx + Inches(0.1), mod_top + Inches(0.28), mod_w - Inches(0.2), mod_h - Inches(0.32))
    set_text(mod_tb.text_frame, mod_title, size=10, color=DARK_TEXT, bold=True)
    add_para(mod_tb.text_frame, mod_desc, size=7.5, color=MEDIUM_TEXT, before=Pt(4))

# ============================
# ROW 3: Key Findings | Customer Segments | Conclusion
# ============================
row3_top = mod_top + mod_h + Inches(0.12)
row3_h = Inches(2.5)

# --- Key Findings ---
find_card = add_rect(slide, c1_left, row3_top, col_w, row3_h, CARD_BG, CARD_BORDER)
find_tb = add_text_box(slide, c1_left + Inches(0.12), row3_top + Inches(0.08), col_w - Inches(0.24), row3_h - Inches(0.12))
section_title(find_tb.text_frame, "Key Findings")

findings = [
    ("Seasonal patterns differ by transaction type", " — Transfers peak mid-year while deposits show end-of-year spikes, confirmed by Kruskal-Wallis tests (p < 0.001)."),
    ("Income brackets drive behavioral divergence", " — High-income customers average 3.2x more transaction volume than low-income segments."),
    ("Geographic clustering of churn risk", " — Rural departments show 40% higher churn probability than urban centers like Bogotá and Medellín."),
    ("Satisfaction score is the #1 churn predictor", " — Lasso feature selection consistently ranks satisfaction_score as the strongest predictor."),
    ("Product bundling reduces churn ~15%", " — What-if simulations show increasing active_products from 1→3 cuts churn from ~47% to ~32%."),
]
for title, desc in findings:
    add_mixed(find_tb.text_frame, [
        ("●  ", TEAL, True),
        (title, DARK_TEXT, True),
        (desc, MEDIUM_TEXT, False),
    ], size=8, before=Pt(3), after=Pt(1))

# --- Discovered Customer Segments ---
seg_card = add_rect(slide, c2_left, row3_top, col_w, row3_h, CARD_BG, CARD_BORDER)
seg_tb = add_text_box(slide, c2_left + Inches(0.12), row3_top + Inches(0.08), col_w - Inches(0.24), Inches(0.5))
section_title(seg_tb.text_frame, "Discovered Customer Segments")
add_para(seg_tb.text_frame, "Optimal clustering (k=4) reveals four distinct behavioral profiles across 10 engineered features:",
         size=8, color=DARK_TEXT, before=Pt(2))

segments = [
    ("Power Users", "~18%", "High frequency, high value, long tenure. Low churn risk. Cross-sell opportunities.", TEAL, TEAL_LIGHT_BG),
    ("Regular Users", "~35%", "Consistent activity, moderate values. Stable segment. Upsell potential via engagement.", BLUE, RGBColor(0xEF,0xF6,0xFF)),
    ("Occasional Users", "~28%", "Sporadic engagement, low product count. Respond well to targeted campaigns.", GREEN, RGBColor(0xEC,0xFD,0xF5)),
    ("At-Risk / New Adopters", "~19%", "Very low activity, short tenure, high churn probability. Priority for onboarding.", RED, RGBColor(0xFE,0xF2,0xF2)),
]

seg_box_top = row3_top + Inches(0.65)
seg_box_h = Inches(0.4)
seg_gap = Inches(0.08)

for i, (name, pct, desc, color, bg_color) in enumerate(segments):
    sy = seg_box_top + i * (seg_box_h + seg_gap)
    # Left accent bar
    accent = add_box(slide, c2_left + Inches(0.15), sy, Pt(4), seg_box_h, color)
    # Segment card
    seg_bg = add_rect(slide, c2_left + Inches(0.19), sy, col_w - Inches(0.36), seg_box_h, bg_color, CARD_BORDER)

    seg_text = add_text_box(slide, c2_left + Inches(0.28), sy + Inches(0.02), col_w - Inches(0.5), seg_box_h)
    set_text(seg_text.text_frame, "", size=1)
    add_mixed(seg_text.text_frame, [
        (name, color, True),
        (f"    {pct} of customers", MUTED_TEXT, False),
    ], size=9, before=Pt(0), after=Pt(0))
    add_para(seg_text.text_frame, desc, size=7, color=MEDIUM_TEXT, before=Pt(1), after=Pt(0))

# --- Conclusion & Recommendations ---
conc_card = add_rect(slide, c3_left, row3_top, col_w, row3_h, CARD_BG, CARD_BORDER)
conc_tb = add_text_box(slide, c3_left + Inches(0.12), row3_top + Inches(0.08), col_w - Inches(0.24), row3_h - Inches(0.12))
section_title(conc_tb.text_frame, "Conclusion & Recommendations")
add_para(conc_tb.text_frame, "COFINFAD Analytics demonstrates that interactive visual analytics can transform raw fintech data into actionable business intelligence.",
         size=8, color=DARK_TEXT, after=Pt(4))

recs = [
    ("1  Retain power users", " — Deploy loyalty programs with real-time satisfaction monitoring dashboards"),
    ("2  Intervene on at-risk segments", " — Personalized outreach and product bundling within 90-day onboarding"),
    ("3  Expand underserved regions", " — Region-specific campaigns for rural Colombian departments with high churn"),
    ("4  Leverage what-if simulations", " — Use the simulator to quantify ROI of retention initiatives before deployment"),
]
for num_title, desc in recs:
    add_mixed(conc_tb.text_frame, [
        (num_title, DARK_TEXT, True),
        (desc, MEDIUM_TEXT, False),
    ], size=8, before=Pt(2), after=Pt(1))

# Recommendation teal number badges
rec_top_base = row3_top + Inches(0.6)
for i in range(4):
    ry = rec_top_base + Inches(i * 0.32)
    rbadge = add_rect(slide, c3_left + Inches(0.14), ry, Inches(0.17), Inches(0.16), TEAL)
    rbadge_tb = add_text_box(slide, c3_left + Inches(0.14), ry - Pt(1), Inches(0.17), Inches(0.16))
    set_text(rbadge_tb.text_frame, str(i+1), size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Future Work
add_para(conc_tb.text_frame, "", size=3, before=Pt(2), after=Pt(0))
add_para(conc_tb.text_frame, "FUTURE WORK", size=9, color=TEAL, bold=True, before=Pt(4), after=Pt(3))
for item in [
    "• Integrate additional data sources (RFM scoring, economic indicators)",
    "• Expand prediction module with ensemble models (RF, XGBoost)",
    "• Add real-time data streaming for live dashboard monitoring",
    "• Incorporate NLP analysis on support ticket text",
]:
    add_para(conc_tb.text_frame, item, size=7.5, color=MEDIUM_TEXT, before=Pt(1), after=Pt(0))

# References
ref_tb = add_text_box(slide, c3_left + Inches(0.12), row3_top + row3_h - Inches(0.5), col_w - Inches(0.24), Inches(0.48))
set_text(ref_tb.text_frame, "", size=1)
add_para(ref_tb.text_frame, "REFERENCES", size=8, color=TEAL, bold=True, before=Pt(0))
for ref in [
    "[1] Kam, T.S. (2024). ISSS608 Visual Analytics & Applications. SMU.",
    "[2] Wickham, H. (2016). ggplot2: Elegant Graphics for Data Analysis.",
    "[3] Patil, I. (2021). ggstatsplot: ggplot2 Based Plots with Statistical Details.",
    "[4] Friedman, J. et al. (2010). Regularization Paths for GLMs via CD. JSS.",
]:
    add_para(ref_tb.text_frame, ref, size=6, color=MUTED_TEXT, before=Pt(0), after=Pt(0))

# ============================
# FOOTER BAR (dark navy)
# ============================
footer_h = Inches(0.35)
footer_top = prs.slide_height - Inches(0.1) - footer_h
footer = add_box(slide, Inches(0.1), footer_top, prs.slide_width - Inches(0.2), footer_h, NAVY_HEADER)
# Teal top accent
add_box(slide, Inches(0.1), footer_top, prs.slide_width - Inches(0.2), Pt(3), TEAL)

footer_left = add_text_box(slide, Inches(0.35), footer_top + Inches(0.08), Inches(8), Inches(0.25))
set_text(footer_left.text_frame, "SMU  •  Singapore Management University  •  ISSS608 Visual Analytics & Applications  •  Term 2, AY2025-26",
         size=8, color=RGBColor(0x8A,0xA3,0xC8))

footer_right = add_text_box(slide, Inches(9), footer_top + Inches(0.08), Inches(4), Inches(0.25))
set_text(footer_right.text_frame, "Team 13: Ji Guofang  •  Lin Yan  •  Nguyen Trong Nhan",
         size=8, color=RGBColor(0x8A,0xA3,0xC8), align=PP_ALIGN.RIGHT)

# ============================
# SAVE
# ============================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "COFINFAD_Poster.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
